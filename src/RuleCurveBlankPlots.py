"""
Create some plots showing the rule curve at all projects on something like the water control diagram (WCD)

Just for testing/playing around

The WCD will start on September 1 so that the winter season remains intact instead of starting on Jan 1
"""

import pandas as pd
import numpy as np
import pickle
from scipy import interpolate
import matplotlib.pyplot as plt
from matplotlib.ticker import MultipleLocator, AutoMinorLocator, FixedLocator, LinearLocator, FormatStrFormatter, PercentFormatter, FuncFormatter, NullFormatter
from scipy.stats import norm
import datetime
import io
from pptx import Presentation
from pptx.util import Inches
from datetime import timedelta

import sys
from pathlib import Path

MODULE_PATH = '../modules'
if not MODULE_PATH in sys.path:
    sys.path.insert(0, MODULE_PATH)

from Utils_Plots import xAxisSeasonal
from ProbabilityDistributions import LP3Dist, ppfLP3, cdfLP3

RULECURVE_CSV_FILE = "../data/common_data/Willamette_Rule_Curves.csv"
DRAWDOWNCURVE_CSV_FILE = "../data/common_data/Drawdown_Target_Curves.csv"
ELEV_STOR_CSV = "../data/common_data/ElevStor.csv" 
RESERVOIR_ELEV_XLSX = "../data/common_data/Reservoirs.xlsx" 
RESIDUAL_LP3_CSV = "../data/LP3_Residual.csv" 


# A list of 11 distinct colors using Tableau 10 and a complementary XKCD color
color_list = [
    'tab:blue',
    'tab:orange',
    'tab:green',
    'tab:red',
    'tab:purple',
    'tab:brown',
    'tab:pink',
    'tab:gray',
    'tab:olive',
    'tab:cyan',
    'xkcd:marigold'  # 11th color: highly distinct gold/yellow
]


def elevToStor(elevInput, resvName):
    '''
    Converts elevation(s) to storage in units of kaf above minimum conservation 
    elevInput can be a number, a list, or an array

    '''
    dfES = dfElevStor.loc[dfElevStor["Reservoir"]==resvName]
    storGross = np.interp(elevInput, dfES["Elev"], dfES["Stor"])
    #Normalize so min conservation pool is 0 storage
    minConElev = resv_info_dict[resvName]["MinConElev"]
    minConStor = np.interp(minConElev, dfES["Elev"], dfES["Stor"])
    storOut = (storGross - minConStor)/1000
    return storOut

def storToElev(storInput, resvName):
    '''
    Reverse of elevToStor, storInput is in kaf above min con

    '''

    dfES = dfElevStor.loc[dfElevStor["Reservoir"]==resvName]
    minConElev = resv_info_dict[resvName]["MinConElev"]
    minConStor = np.interp(minConElev, dfES["Elev"], dfES["Stor"])
    storGross = storInput*1000 + minConStor
    elevOut = np.interp(storGross, dfES["Stor"], dfES["Elev"])
    return elevOut

def getIRRMLine(irrmElev, resvName, seriesRCStor):
    '''
    Given a series of rule curve storage,
    Return a series with the IRRM as a flat line (only portions below the rule curve)
    This is for seismic IRRM
    '''
    irrmStor = elevToStor(irrmElev, resvName)
    seriesIRRM = seriesRCStor.copy()
    seriesIRRM.loc[seriesRCStor.values < irrmStor] = np.nan
    seriesIRRM.loc[seriesRCStor.values >= irrmStor] = irrmStor
    #seriesIRRM.loc[:] = irrmStor
    return seriesIRRM

def yAxisStorAtKeyElevs(ax, resvName):
    '''
    Formats the y-axis the way we want it for stor
    If showElev=True, the secondary y-axis will show elevation
    '''
    keyElevColNames = ["MinElev","FullPoolElev","MinConElev","MaxConElev","MinPowerPool","InactivePool","IRRMElev","SpillwayCrest"]
    elevs = sorted(dfReservoirs.loc[dfReservoirs["ReservoirName"] == resvName, keyElevColNames].squeeze().dropna())
    stors = elevToStor(elevs, resvName)
    ax.yaxis.set_major_locator(FixedLocator(stors))
    ax.yaxis.set_major_formatter(FormatStrFormatter('%d'))
    #Add the elevations at right
    addSecondaryYLabelsElev(ax, resvName)

def addSecondaryYLabelsElev(ax, resvName):
    ax2 = ax.twinx()
    #ax2.spines["top"].set_visible(False)
    #ytick2 = ((fullPoolStor - interpElevStor(ytick))/1000).round(2)
    ax2.set_ylim(ax.get_ylim())
    ticksLeft = ax.get_yticks()
    elevs = storToElev(ticksLeft, resvName) 
    yRightLabels = ["%0.1f" %elev for elev in elevs]
    ax2.yaxis.set_major_locator(FixedLocator(ticksLeft))
    ax2.set_yticklabels(yRightLabels)
    ax2.set_ylabel("Elevation (ft, project datum)")

dfRuleCurves = pd.read_csv(RULECURVE_CSV_FILE, parse_dates=["Date"], index_col="Date")
dfDrawdownCurves = pd.read_csv(DRAWDOWNCURVE_CSV_FILE, parse_dates=["Date Time"], index_col="Date Time")
dfElevStor = pd.read_csv(ELEV_STOR_CSV)
dfReservoirs = pd.read_excel(RESERVOIR_ELEV_XLSX, sheet_name="Reservoirs", header = 5)
resv_info_dict = dfReservoirs.set_index('ReservoirName').to_dict(orient='index') #can call resv_info_dict["Detroit"]["SpillwayCrest"]
dfLP3Residual = pd.read_csv(RESIDUAL_LP3_CSV, parse_dates=["Date"])

startDT = datetime.datetime(2000,8,1)
endDT = datetime.datetime(2001,7,31)

dfOneYear = dfRuleCurves.loc[(dfRuleCurves.index >= startDT) & 
                             (dfRuleCurves.index <= endDT)]

#Convert elevation to storage in kaf above min con
dfRCStor = dfOneYear.copy()
for resvName in dfOneYear.columns:
    dfRCStor[resvName] = elevToStor(dfOneYear[resvName], resvName)
    
dfRCPct = dfRCStor/dfRCStor.max()

plt.close('all')
#Spaghetti plot of rule curve storage
fig, ax = plt.subplots()
for i, resvName in enumerate(dfRCStor.max().sort_values(ascending=False).index):
    ax.plot(dfRCStor[resvName], label=resvName, color = color_list[i])
    #Add text label
    ax.text(dfRCStor.index[-1], dfRCStor[resvName].max(), f" {resvName}", ha="left", va="center")
ax.set_xlim([dfRCStor.index[0],dfRCStor.index[-1]])
ax = xAxisSeasonal(ax)
ax.set_ylabel("Rule Curve Storage (kaf)")
plt.tight_layout()


#Same spaghetti plot, but with percent of rule curve storage
fig, ax = plt.subplots()
#for resvName in dfRCPct.loc[datetime.datetime(2001,3,1)].sort_values(ascending=False).index:
#Keep same colors as before
for i, resvName in enumerate(dfRCStor.max().sort_values(ascending=False).index):
    ax.plot(dfRCPct[resvName], label=resvName, color = color_list[i])
    #Add text label
    #ax.text(dfRCPct.index[-1], dfRCPct[resvName].max(), f" {resvName}", ha="left", va="center")
ax.set_xlim([dfRCPct.index[0],dfRCPct.index[-1]])
plt.legend()
ax = xAxisSeasonal(ax)
ax.yaxis.set_major_formatter(PercentFormatter(xmax=1))
ax.set_ylabel("Percent of Rule Curve Storage")
plt.tight_layout()


#Plot up the rule curve as more of a simplified cartoon with zones colored differently

# Initialize presentation
prs = Presentation()
# Configure 16:9 Widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use layout 6 (Blank slide) or 5 (Title only)
blank_layout = prs.slide_layouts[6]

for resvName in dfReservoirs["ReservoirName"]:
#for resvName in ["Lookout Point"]:
    fig, ax = plt.subplots(figsize=[8, 4.5])
    ax.plot(dfRCStor[resvName], color='black')
    ax = xAxisSeasonal(ax)
    ax.set_xlim([dfRCStor.index[0], dfRCStor.index[-1]])
    ax.set_ylabel("Storage (kaf)")
    fig.suptitle(resvName)
    #Add a bunch of markers
    
    if resvName != "Foster":
        spillCrestElev = resv_info_dict[resvName]["SpillwayCrest"]
        spillCrestStor = elevToStor(spillCrestElev, resvName)
        ax.axhline(spillCrestStor, ls="dashed", color="gray", alpha=0.3, zorder=2)
        ax.text(dfRCStor.index[-2], elevToStor(spillCrestElev-0.5, resvName), "Spillway Crest", ha="right",va="top")
    
    #Shade the various pools
    
    #Flood pool
    fullStor = elevToStor(resv_info_dict[resvName]["FullPoolElev"], resvName)
    ax.fill_between(
        dfRCStor[resvName].index, 
        fullStor, 
        dfRCStor[resvName], 
        color='tab:blue',             # Selected soft color from the "Cool" palette
        alpha=0.4,                   # Transparency so background grid is visible
        zorder=1                     # Places the fill behind the boundary lines
    )
    ax.text(datetime.datetime(2001,1,1), 0.3*dfRCStor[resvName].max(), "Flood Pool", ha="center",va="center")
    
    #Seismic IRRM Pool
    seriesRCClipped = dfRCStor[resvName].copy()
    irrmElev = resv_info_dict[resvName]["IRRMElev"]
    if not pd.isna(irrmElev):
        irrmStor = elevToStor(irrmElev, resvName)
        seriesIRRM = getIRRMLine(irrmElev, resvName, dfRCStor[resvName])
        #ax.plot(seriesIRRM, ls="dashed", color="tab:red", zorder = 2)
        ax.fill_between(
            seriesIRRM.index, 
            seriesIRRM.values, 
            dfRCStor[resvName], 
            where=(dfRCStor[resvName] >= seriesIRRM.values),      # Only shade where line1 is above line2
            interpolate=True,            # Ensures a clean fill at the intersection points
            color='tab:red',             # Selected soft color from the "Cool" palette
            alpha=0.4,                   # Transparency so background grid is visible
            zorder=1                     # Places the fill behind the boundary lines
        )
        seriesRCClipped = seriesRCClipped.clip(upper=irrmStor)
        ax.text(dfRCStor.index[-2], (dfRCStor[resvName].max()+irrmStor)/2, "Seismic IRRM Pool", ha="right",va="center")
    
    #Conservation Pool
    lowStor = 0
    ax.fill_between(
        dfRCStor[resvName].index, 
        lowStor, 
        seriesRCClipped, 
        interpolate=True,            # Ensures a clean fill at the intersection points
        color='tab:green',             # Selected soft color from the "Cool" palette
        alpha=0.4,                   # Transparency so background grid is visible
        zorder=1                     # Places the fill behind the boundary lines
    )
    ax.text(datetime.datetime(2001,4,15), 0.3*dfRCStor[resvName].max(), "Conservation Pool", ha="center", va="center")
    
    #Min power pool
    minPowerPoolElev = resv_info_dict[resvName]["MinPowerPool"]
    if not pd.isna(minPowerPoolElev):
        minPowerPoolStor = elevToStor(minPowerPoolElev, resvName)
        ax.axhspan(
            lowStor, 
            minPowerPoolStor, 
            color='tab:purple',             # Selected soft color from the "Cool" palette
            alpha=0.4,                   # Transparency so background grid is visible
            zorder=1                     # Places the fill behind the boundary lines
        )
        lowStor = minPowerPoolStor
        ax.text(datetime.datetime(2001,4,15), 0.5*minPowerPoolStor, "Min Power Pool", ha="center", va="center")
    
    #Drawdown pool
    if resvName in dfDrawdownCurves.columns:
        dfDD = dfDrawdownCurves.loc[(dfDrawdownCurves.index >= dfRCStor.index[0]) & 
                                    (dfDrawdownCurves.index <= dfRCStor.index[-1])]
        seriesDD = dfDD[resvName].dropna()
        seriesDD[:] = elevToStor(seriesDD, resvName)
        ax.plot(seriesDD, color="tab:brown", ls = 'dotted', lw=2)
        if resvName == "Lookout Point":
            ax.text(datetime.datetime(2000,10,2), 145, "Deep drawdown", va='top',ha='left',rotation=-73)
            ax.plot([datetime.datetime(2001,3,12),
                     datetime.datetime(2001,4,12),
                     datetime.datetime(2001,4,13)],
                    [spillCrestStor+4, spillCrestStor+4, elevToStor(908, resvName)],
                    color="tab:brown", ls = 'dotted', lw=2)
            ax.text(datetime.datetime(2001,3,13), spillCrestStor+2, "Spill Op", va='top',ha='left',rotation=0)
        if resvName == "Green Peter":
            ax.text(datetime.datetime(2000,10,2), 130, "Deep drawdown", va='top',ha='left',rotation=-73)
        if resvName == "Fall Creek":
            ax.text(datetime.datetime(2000,10,16), 30, "Deep drawdown", va='top',ha='left',rotation=-78)
        if resvName == "Detroit":
            ax.text(datetime.datetime(2000,10,16), 80, "Deep drawdown", va='top',ha='left',rotation=-67)
        if resvName == "Cougar":
            ax.text(datetime.datetime(2000,10,5), 40, "Deep drawdown", va='top',ha='left',rotation=-76)
        if resvName == "Foster":
            ax.text(datetime.datetime(2001,2,10), 2, "Delayed Refill", va='top',ha='left',rotation=0)
    drawdownPoolElev = resv_info_dict[resvName]["DrawdownPool"]
    #if not pd.isna(drawdownPoolElev):
    if 1 == 0:
        drawdownPoolStor = elevToStor(drawdownPoolElev, resvName)
        ax.axhspan(
            lowStor, 
            drawdownPoolStor, 
            color='tab:brown',             # Selected soft color from the "Cool" palette
            alpha=0.4,                   # Transparency so background grid is visible
            zorder=1                     # Places the fill behind the boundary lines
        )
        lowStor = drawdownPoolStor
    
    #Inactive Pool
    #Use "inactive" if defined, otherwise use "min pool", which is lowest outlet invert
    inactivePoolElev = resv_info_dict[resvName]["InactivePool"]
    minPoolElev = resv_info_dict[resvName]["MinElev"]
    minPoolStor = elevToStor(minPoolElev, resvName)
    if not pd.isna(inactivePoolElev):
        minPoolStor = elevToStor(inactivePoolElev, resvName)
        
    ax.axhspan(
        lowStor, 
        minPoolStor, 
        color='tab:orange',             # Selected soft color from the "Cool" palette
        alpha=0.4,                   # Transparency so background grid is visible
        zorder=1                     # Places the fill behind the boundary lines
    )
    ax.text(datetime.datetime(2001,4,15), (lowStor + minPoolStor)/2, "Inactive Pool", ha="center", va="center")
    
    ax.set_ylim([minPoolStor, fullStor])
    ax.grid(True, axis='x', alpha=0.5)
    yAxisStorAtKeyElevs(ax, resvName)
    plt.tight_layout()
    
    
    #Push to powerpoint
    # Save to memory buffer
    image_stream = io.BytesIO()
    fig.savefig(image_stream, format="png", bbox_inches="tight")
    image_stream.seek(0)
    #plt.close(fig)  # Free memory

    # Add slide and insert image
    slide = prs.slides.add_slide(blank_layout)
    
    # Position: left, top, width (height auto-scales to maintain aspect ratio)
    left = Inches(2.66)
    top = Inches(1.5)
    width = Inches(8.0)
    
    #slide.shapes.add_picture(image_stream, left, top, width=width)
    
    slide.shapes.add_picture(
        image_stream,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height
    )

# Save final deck
prs.save("../out/reservoir_plots.pptx")
    

# Plot up the rule curve space vs the probability of flood it takes to fill it in 3 days
#Kind of like the risk rule curve
fig, ax = plt.subplots(figsize=[8, 4.5])
#for resvName in ["Detroit"]:
for resvName in dfReservoirs["ReservoirName"]:
    if resvName =="Foster": continue #meaningless
    #Get the n-day flow required to fill flood space for each rule curve ordinate
    durationDays = 5
    fullStor = elevToStor(resv_info_dict[resvName]["FullPoolElev"], resvName)
    series3DayFlow = (fullStor - dfRCStor[resvName])*43560/(durationDays*24*60*60)*1000 #average cfs over 3 days
    #Deduct out min flow
    series3DayFlow = series3DayFlow + resv_info_dict[resvName]["MinFlowFlood"]
    if resvName =="Lookout Point":
        #Add Hills Creek storage too
        fullStorHCR = elevToStor(resv_info_dict["Hills Creek"]["FullPoolElev"], "Hills Creek")
        series3DayFlowHCR = (fullStorHCR - dfRCStor["Hills Creek"])*43560/(durationDays*24*60*60)*1000 #average cfs over 3 days    
        series3DayFlow = series3DayFlow + series3DayFlowHCR
    #Get the residual probability distributions
    dfLP3 = dfLP3Residual.loc[(dfLP3Residual["LocationName"]==resvName) &
                               (dfLP3Residual["Duration"]=="5-Day")]
    #Get the probability
    rows = []
    for dt, flow3Day in series3DayFlow.items():
        dt2000 = datetime.datetime(2000, dt.month, dt.day)
        rowLP3 = dfLP3.loc[dfLP3["Date"]==dt2000]
        if len(rowLP3) == 0: continue
        #Don't show information after rule curve has reached full
        if dt.month >= 4:
            if flow3Day == series3DayFlow[dt-timedelta(days=1)]:
                continue
        dictLP3 = rowLP3.to_dict('records')[0]
        distLP3 = LP3Dist(dictLP3["Mean"], dictLP3["StdDev"], dictLP3["Skew"])
        aep = 1 - cdfLP3(distLP3, flow3Day)
        zScore = -norm.ppf(aep)
        rows.append([dt, aep])
    dfPlot = pd.DataFrame(data=rows, columns=["Date","AEP"]).set_index("Date")
    ax.plot(dfPlot.index, -norm.ppf(dfPlot["AEP"]), label=resvName)
ax.yaxis.set_major_formatter(PercentFormatter(xmax=1))
ax.set_ylabel("Chance of a 5-day storm volume large enough to fill the flood pool")


#Primary y-axis labels - Percent Annual Exceedance Probability
ax.set_axisbelow(True) #makes the grid lines show up behind everything else
ax.grid() #adds grid
prob = np.array([0.90, 0.50, 0.10, 0.01, 0.001, 0.0001, 0.00001])
probStr = ["90%","50%","10%","1%","0.1%","0.01%", "0.001%"]
#prob_percent = np.multiply(prob,100)
normprob_ticks = -norm.ppf(prob)  #calculates the inverse of standard normal cumulative distribution function for desired axis labels
#plt.xticks(normprob_ticks,prob_percent.astype(str)) # set tick locations
#plt.xticks(normprob_ticks,probStr) # set tick locations
ax.set_yticks(normprob_ticks,probStr) # set tick locations
# define the minor grid lines 
prob_minor = np.array([0.90, 0.80, 0.70, 0.60, 0.50, 0.40, 0.30, 0.20, 0.10, 
    0.09, 0.08, 0.07, 0.06, 0.05, 0.04, 0.03, 0.02, 0.01, 
    0.009, 0.008, 0.007, 0.006, 0.005, 0.004, 0.003, 0.002, 0.001,
    0.0009, 0.0008, 0.0007, 0.0006, 0.0005, 0.0004, 0.0003, 0.0002, 0.0001,
    0.00009, 0.00008, 0.00007, 0.00006, 0.00005, 0.00004, 0.00003, 0.00002, 0.00001])
normprob_minticks = -norm.ppf(prob_minor)
#Show minor gridlines with very faint and almost transparent grey lines
ax.yaxis.set_minor_locator(FixedLocator(normprob_minticks))
ax.grid(visible=True, which='minor', axis = 'y',color='#999999', linestyle='-', alpha=0.2)
#Set limits based on specified decimal high and low annual exceedance probability
lowprob=0.001
highprob=0.99
ax.set_ylim([-norm.ppf(lowprob),-norm.ppf(highprob)])


ax = xAxisSeasonal(ax)
ax.set_xlim([dfPlot.index[0], datetime.datetime(2001,6,1)])
ax.legend()


